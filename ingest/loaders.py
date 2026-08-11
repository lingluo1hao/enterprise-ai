"""多格式文档加载器（改造点 3）。

统一从「仅 .pdf/.txt」扩展到 docx/xlsx/pptx/html/md，并保留表格结构：
- .txt/.md   → TextLoader（零额外依赖）
- .pdf       → PyPDFLoader（项目已有）
- .html/.htm → BSHTMLLoader（带 bs4）；缺失则降级 TextLoader
- .docx      → Docx2txtLoader；缺失则降级 UnstructuredWordDocumentLoader
- .xlsx/.xls → openpyxl 逐 Sheet 逐行读（保留表头/单元格，利于表格问答）
- .pptx      → python-pptx 逐页读文本+表格

所有 loader 均为「懒导入」：缺失某个可选依赖时只跳过该格式，不拖垮整个模块，
保证 ingest 引擎在无全套解析依赖的机器上仍能跑核心链路（txt/md/pdf）。
"""

import glob
import os
import re
from typing import List, Optional, Any

import statistics

from .types import RawDoc
from .chunk import FIGURE_ANCHOR_RE

# 默认权限规则（与 advanced_rag_agent.DOC_ACCESS_RULES 语义一致）。
# 生产路径下由 VectorStoreManager 传入 AccessControlFilter.get_access_level 覆盖。
# 规则外置到 access_rules.yaml：filename 关键字 -> 级别（public/restricted）；
# 缺失文件或 yaml 不可用则退化为空（默认全部 public）。
DOC_ACCESS_RULES: dict = {}
try:
    import yaml  # 可选依赖：未安装时退化为全 public
    _cfg_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                              "..", "config", "access_rules.yaml")
    if os.path.isfile(_cfg_path):
        with open(_cfg_path, "r", encoding="utf-8") as _fh:
            _cfg = yaml.safe_load(_fh) or {}
        DOC_ACCESS_RULES = {str(k): v for k, v in (_cfg.get("access_rules") or {}).items()}
except Exception:
    DOC_ACCESS_RULES = {}


def get_access_level(source: str) -> str:
    """按文件名关键字判定权限：命中 → 对应级别，否则 public。"""
    basename = os.path.basename(source)
    for keyword, level in DOC_ACCESS_RULES.items():
        if keyword in basename:
            return level
    return "public"


def _raw(text: str, source: str) -> List["object"]:
    from .types import RawDoc
    return [RawDoc(text=text, source=source, file_name=os.path.basename(source))]


def _load_xlsx(path: str) -> Optional[list]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        lines = [f"# Sheet: {ws.title}"]
        for r, row in enumerate(ws.iter_rows(values_only=True), 1):
            cells = [str(c).strip() for c in row
                     if c is not None and str(c).strip() != ""]
            if cells:
                lines.append(f"R{r}: " + " | ".join(cells))
        if len(lines) > 1:
            sheets.append("\n".join(lines))
    if not sheets:
        return None
    return _raw("\n\n".join(sheets), path)


def _load_pptx(path: str) -> Optional[list]:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip()
                                            for c in row.cells if c.text.strip()))
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    if not slides:
        return None
    return _raw("\n\n".join(slides), path)


def _table_to_markdown(rows: List[List[Any]]) -> Optional[str]:
    """把 PyMuPDF 表格二维数组转成标准 Markdown 表格。"""
    if not rows or len(rows) < 1:
        return None
    # 清理单元格：去空值、压平换行、转义 |
    def _cell(v):
        if v is None:
            return ""
        s = str(v).replace("\n", " ").replace("\r", " ").replace("|", "\\|")
        return s.strip()
    cleaned = [[_cell(c) for c in row] for row in rows]
    n_cols = max(len(r) for r in cleaned)
    # 补齐列数，避免 Markdown 格式错乱
    for r in cleaned:
        while len(r) < n_cols:
            r.append("")
    lines = []
    for i, row in enumerate(cleaned):
        lines.append("| " + " | ".join(row) + " |")
        if i == 0:
            lines.append("| " + " | ".join(["---"] * n_cols) + " |")
    return "\n".join(lines)


def _figures_dir(path: str) -> tuple:
    """返回 PDF 图/表输出目录的 (绝对路径, 相对项目根目录的 URL 前缀)。"""
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if not os.path.isdir(os.path.join(project_root, "knowledge")):
        project_root = os.getcwd()
    stem = os.path.splitext(os.path.basename(path))[0]
    out_dir = os.path.join(project_root, "assets", "figures", stem)
    rel_prefix = f"assets/figures/{stem}"
    return out_dir, rel_prefix


def _probe_font_levels(pages_items):
    """字体大小兜底探测（方案丙）。

    当文档没有数字编号标题（split_level 为 None）时，用每文本块的最大字号
    识别两级标题：比正文基字号大一个明显台阶的当作标题，最大档为 lvl2（章节）、
    次大档为 lvl1（小节）。返回 (font_mode: bool, levels: dict)。
    失败或缺数据均降级为 (False, {})，不影响主流程。
    """
    sizes = [tb.get("max_font", 0.0) for _, items, _ in pages_items
             for _, kind, tb in items
             if kind == "text" and tb.get("max_font", 0.0) > 0]
    if not sizes:
        return False, {}
    try:
        body = statistics.median(sizes)
    except statistics.StatisticsError:
        return False, {}
    THRESH = 1.15   # 字号大于正文 15% 才算候选标题
    cand = []
    for page_idx, items, _ in pages_items:
        for y0, kind, tb in items:
            if kind != "text":
                continue
            ms = tb.get("max_font", 0.0)
            txt = tb.get("text", "")
            if ms > body * THRESH and 0 < len(txt) <= 80:
                cand.append((page_idx, y0, txt, ms))
    if not cand:
        return False, {}
    cand_sizes = sorted({ms for *_, ms in cand})
    top2 = cand_sizes[-2:]
    levels = {}
    for page_idx, y0, txt, ms in cand:
        lvl = 2 if (len(top2) >= 2 and ms >= top2[-1]) else 1
        levels[(page_idx, y0, txt)] = lvl
    return True, levels


def _is_valid_table(rows, bbox, page_rect) -> bool:
    """结构性校验：拒绝 PyMuPDF 把文字段误检成的伪表格。

    阈值（实测学生证 PDF 校准，见修复报告）：
      ① rows >= 2           —— 单行"表"= 文字段被误框（p8/p14/p45 等误检）
      ② cols >= 2           —— 单列"表"= 列表被误框
      ③ max_cell_len <= 150 —— 整段叙述被当 1 个 cell，合并误检主因（用户截图 p2=185）
      ⑤ bbox 占比 <= 0.72   —— 整页文字被当一张表（实测最大 0.67）
    命中任一拒绝条件即判为伪表格。数据不丢：调用方对伪表格仍保留 [TABLE] 文本块，
    仅跳过 PNG 渲染与 figure_paths 关联。
    """
    if not rows:
        return False
    ncols = max((len(r) for r in rows), default=0)
    if len(rows) < 2 or ncols < 2:
        return False
    cells = [c for r in rows for c in r if c is not None]
    if not cells:
        return False
    lens = [len(str(c).strip()) for c in cells]
    if max(lens) > 150:
        return False
    try:
        ratio = ((bbox[2] - bbox[0]) * (bbox[3] - bbox[1])
                 / (page_rect.width * page_rect.height))
    except Exception:
        ratio = 0.0
    if ratio > 0.72:
        return False
    return True


def _load_pdf(path: str) -> Optional[list]:
    """使用 PyMuPDF 提取 PDF 文本与表格，并**按章节聚合为结构化单元**。

    核心改造（解决「按页切 → 章节破碎、图/表/文错配」）：
    - 全文档扫描标题层级（如 2 / 2.13 / 2.13.2），选定切分级别（默认二级 2.13）；
    - 以「章节」为聚合单元：每个章节 = 其标题 + 下属文字 + 下属表格（[TABLE]+[[FIG:]]）
      + 该章节范围内的所有图片（通用图抽取 + 表格图）。**大章节天然包含小章节**
      （小章节标题作为普通文本内嵌，表格/图片/文字都跟随所属章节）；
    - 每个章节产出一个 RawDoc（section_path 已设置），交由 chunker 做
      「整章为父、细切为子」的 small-to-big 父子切片；
    - 表格区域渲染为 PNG 并内嵌 [[FIG:...]]，表格图不挂到页级 figure_paths。
    - 缺失 PyMuPDF 或提取失败时返回 None，由调用方降级到 PyPDFLoader。
    """
    try:
        import fitz
    except ImportError:
        return None
    try:
        out_dir, rel_prefix = _figures_dir(path)
        os.makedirs(out_dir, exist_ok=True)
        # 清理旧版表格图，避免重建后残留过期图片
        for old in glob.glob(os.path.join(out_dir, "table_p*.png")):
            try:
                os.remove(old)
            except OSError:
                pass

        # 通用图抽取（每页列表，与语言/caption/排版无关）
        gen_figs = _extract_figures(path)  # 失败返回 []

        # 标题识别：编号标题（2 / 2.13 / 2.13.2 ...），过长的行不当标题
        def _heading_level(t: str):
            t = t.strip()
            if len(t) > 80:
                return None
            # 数字编号标题：2 / 2.13 / 2.13.2
            m = re.match(r"^(\d+(?:\.\d+)+)\s+", t)
            if m:
                return m.group(1).count(".") + 1  # 2→1, 2.13→2, 2.13.2→3
            # 罗马数字编号章节：VI 基站信息格式 / IV 技术要求 ...
            ROMANS = ['X', 'IX', 'VIII', 'VII', 'VI', 'V', 'IV', 'III', 'II', 'I']
            for r in ROMANS:
                if re.match(rf"^{r}\s+[\u4e00-\u9fff]", t):
                    return 2
            # 中文编号小节：一、 二、 三．  （学生证协议等常用）
            if re.match(r'^[一二三四五六七八九十百]+\s*[、.．]', t):
                return 2
            # 数字加顿号/点小节：1、一般信息 / 1. 概述（ASCII 小数点 1.0 不算标题）
            if re.match(r'^\d+\s*[、．]', t) or re.match(r'^\d+\.(?!\d)\s', t):
                # 排除「1、2、3 号键拨出的号码…」这类含多个顿号/句号的列表句
                if t.count("、") <= 1 and "。" not in t:
                    return 1
            # 裸中文短标题（无编号前缀）：协议类文档常见如「基站信息格式」「关于白名单」。
            # 约束：3~14 个纯汉字、不含句末标点，且以标题后缀（格式/结构/信息/协议/命令/
            # 流程/集/说明/规则/要求/规范/表/图）结尾，或以「关于/附录/附件」开头。
            # 这样既能把纯中文标题切成独立章节，又不会把正文短句（设备登陆/其中：…）
            # 或表格脚注（裸「格式」/「说明」）误判成标题。
            TITLE_SUFFIX = ('格式', '结构', '信息', '协议', '命令', '流程', '集',
                            '说明', '规则', '要求', '规范', '表', '图')
            if (3 <= len(t) <= 14
                    and re.match(r'^[\u4e00-\u9fff]+$', t)
                    and (t.endswith(TITLE_SUFFIX)
                         or t.startswith(('关于', '附录', '附件')))):
                return 2
            return None

        # 第一遍：按页提取 items（text / table），并统计标题层级分布
        pages_items = []          # [(page_idx, [(kind, payload), ...], [gen_fig_paths])]
        level_counts = {}
        with fitz.open(path) as doc:
            for page_idx, page in enumerate(doc, 1):
                # 1) 表格
                tables = []
                try:
                    tabs = page.find_tables()
                    for tidx, tab in enumerate(tabs.tables if tabs else []):
                        try:
                            rows = tab.extract()
                        except Exception:
                            continue
                        md = _table_to_markdown(rows)
                        if not md:
                            continue
                        bbox = fitz.Rect(tab.bbox)
                        # 结构性校验：拒绝伪表格（整段文字被误框成表）。
                        # 保留 [TABLE] 文本块（不丢数据），仅跳过 PNG 渲染与 figure_paths。
                        if not _is_valid_table(rows, bbox, page.rect):
                            tables.append({"y0": bbox.y0, "bbox": bbox, "md": md, "fig_path": ""})
                            continue
                        fig_name = f"table_p{page_idx:03d}_{tidx+1}.png"
                        fig_path = f"{rel_prefix}/{fig_name}"
                        try:
                            scale = 2.0
                            pad = 6
                            clip = fitz.Rect(
                                max(0, bbox.x0 - pad), max(0, bbox.y0 - pad),
                                min(page.rect.width, bbox.x1 + pad),
                                min(page.rect.height, bbox.y1 + pad),
                            )
                            pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), clip=clip)
                            pix.save(os.path.join(out_dir, fig_name))
                        except Exception as te:
                            print(f"[ingest] 第 {page_idx} 页表格 {tidx+1} 渲染失败: {te}")
                            fig_path = ""
                        tables.append({"y0": bbox.y0, "bbox": bbox, "md": md, "fig_path": fig_path})
                except Exception as te:
                    print(f"[ingest] 第 {page_idx} 页表格检测失败: {te}")

                # 2) 非表格文本块（顺带取块内最大字号，供字体标题兜底探测）
                text_blocks = page.get_text("blocks")
                raw_dict = page.get_text("dict")
                dict_blocks = raw_dict.get("blocks", [])  # 保留全量，与 text_blocks 按序对齐
                text_items = []
                for bi, b in enumerate(text_blocks):
                    bbox = fitz.Rect(b[:4])
                    txt = b[4].strip()
                    if not txt:
                        continue
                    # 文字块要不要跳过，**必须看文字本身**是不是这个表的 cell 内容。
                    # 之前用 "bbox 相交且 md 非空就跳过" 在 yh 协议 PDF 第 8 页误伤：fitz find_tables()
                    # 把"命令集 S23 服务器指令"识别成跨整页的大表 bbox(84.6,188.8 → 728.0,837.9)，
                    # 表 md 还是有内容（1065 字符），但 bbox 越过表实际范围、把下面的「2.2 工作模式设置：
                    # MODE / 产品有4 种工作模式…」整节正文也圈了进去，按旧逻辑会被当表内文字吞掉。
                    # 这里改为：文字块 bbox 与某表相交 **且** 该文字确实出现在该表的 cell 文本/md 中，才跳过；
                    # 否则一律保留正文（真表内 cell 文字几乎都出现在 md 里，绝不会误伤）。
                    _skip_block = False
                    for t in tables:
                        if not bbox.intersects(t["bbox"]):
                            continue
                        md_text = t.get("md") or ""
                        if not md_text.strip():
                            continue
                        # 短 cell 直接包含（如 "格式"）；长 cell 用首 16 字符前缀匹配（避免 memory hits）
                        if txt in md_text:
                            _skip_block = True
                            break
                    if _skip_block:
                        continue
                    max_font = 0.0
                    if bi < len(dict_blocks):
                        for ln in dict_blocks[bi].get("lines", []):
                            for sp in ln.get("spans", []):
                                max_font = max(max_font, float(sp.get("size", 0.0)))
                    text_items.append({"y0": bbox.y0, "text": txt, "max_font": max_font})

                # 3) 按阅读顺序合并，并登记标题层级
                items = []
                for t in tables:
                    items.append((t["y0"], "table", t))
                for tb in text_items:
                    items.append((tb["y0"], "text", tb))
                items.sort(key=lambda x: x[0])

                for _, kind, blk in items:
                    if kind == "text":
                        lv = _heading_level(blk["text"])
                        if lv is not None:
                            level_counts[lv] = level_counts.get(lv, 0) + 1

                gfigs = gen_figs[page_idx - 1] if page_idx - 1 < len(gen_figs) else []
                pages_items.append((page_idx, items, gfigs))

        # 选定切分级别：优先二级（2.13），否则一级（2），否则整文档一章
        if level_counts.get(2, 0) >= 1:
            split_level = 2
        elif level_counts.get(1, 0) >= 1:
            split_level = 1
        else:
            split_level = None

        # 双轨兜底（方案丙）：数字编号标题为零时，用字号大小识别两级标题
        font_levels: dict = {}     # (page_idx, y0, text) -> lvl(1|2)
        font_mode = False
        if split_level is None:
            font_mode, font_levels = _probe_font_levels(pages_items)
            if font_mode:
                lv2 = sum(1 for v in font_levels.values() if v == 2)
                split_level = 2 if lv2 >= 1 else 1

        # 第二遍：按 split_level 标题切分章节，聚合文字/表格/图片。
        # 单章节超限（> MAX_SECTION_CHARS）自动「续切」为同名续章，
        # 既保证「整章为父窗口不超 Milvus 字段上限」，又让大章节仍能完整还原。
        MAX_SECTION_CHARS = 7800
        sections = []
        current = None
        chapter_prefix = ""   # 最近的上级章节标题（如 "2 WIFI 定位"），用于层级路径

        def _new_section(title="", start_page=None, path=None):
            return {"items": [], "figures": [], "title": title,
                    "start_page": start_page, "path": path or [], "length": 0}

        def _has_content(s):
            for kind, _ in s["items"]:
                if kind in ("text", "table", "fig"):
                    return True
            return False

        def _render(s):
            parts = []
            for kind, val in s["items"]:
                if kind == "text":
                    parts.append(val)
                elif kind == "fig":
                    parts.append(f"[[FIG:{val}]]")
                elif kind == "table":
                    fig_line = f"[[FIG:{val['fig_path']}]]\n" if val["fig_path"] else ""
                    parts.append(f"[TABLE]\n{fig_line}{val['md']}\n[/TABLE]")
            return "\n\n".join(p for p in parts if p.strip())

        def _flush():
            nonlocal current
            if current is not None and _has_content(current):
                sections.append({
                    "text": _render(current),
                    "figures": current["figures"],
                    "path": current["path"],
                    "start_page": current["start_page"],
                })

        def _add_item(kind, val, length):
            nonlocal current
            # 续切：当前章节已有内容且加上本项将超限 → 收尾并开同名续章（路径不变）
            if (current is not None and _has_content(current)
                    and current["length"] + length > MAX_SECTION_CHARS):
                _flush()
                current = _new_section(title=current["title"],
                                       start_page=current["start_page"],
                                       path=list(current["path"]))
            if current is None:
                current = _new_section(start_page=None)
            current["items"].append((kind, val))
            current["length"] += length
            if kind == "table" and val["fig_path"]:
                current["figures"].append(val["fig_path"])
            elif kind == "fig":
                current["figures"].append(val)

        for page_idx, items, gfigs in pages_items:
            for _, kind, blk in items:
                if kind == "table":
                    _add_item("table", blk, len(blk["md"]) + 40)
                else:  # text
                    text = blk["text"]
                    lv = _heading_level(text)
                    if font_mode and lv is None:
                        lv = font_levels.get((page_idx, blk["y0"], text))
                    if split_level is not None and lv == split_level:
                        # 章节边界：收尾上一章，开新章，并把标题行写入新章
                        _flush()
                        current = _new_section(title=text, start_page=page_idx)
                        current["path"] = ([chapter_prefix] if chapter_prefix else []) + [text]
                        _add_item("text", text, len(text))
                    else:
                        # 非切分级标题：更新上级章节前缀（仅一级），并作为正文行内嵌
                        if lv == 1 and split_level == 2:
                            chapter_prefix = text
                        _add_item("text", text, len(text))
            # 本页通用图：挂到当前章节，并按出现位置内嵌 [[FIG:]]
            for fg in gfigs:
                _add_item("fig", fg, 40)
        _flush()

        # 章节 → RawDoc（section_path 已设置，标记为「已分章节单元」）
        out = []
        for sec in sections:
            figs = list(dict.fromkeys(sec["figures"]))  # 去重保序
            out.append(RawDoc(
                text=sec["text"],
                source=path,
                file_name=os.path.basename(path),
                page=sec["start_page"],
                figure_paths=figs,
                section_path=sec["path"] or None,
            ))

        if not out:
            return None
        return out
    except Exception as e:
        print(f"[ingest] PyMuPDF 读取 PDF 失败 {os.path.basename(path)}: {e}")
        return None


def load_file(path: str) -> List["object"]:
    """加载单个文件，返回 RawDoc 列表。不支持/缺依赖的格式返回 []（打印警告）。"""
    from .types import RawDoc
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in (".txt", ".md"):
            from langchain_community.document_loaders import TextLoader
            docs = TextLoader(path, encoding="utf-8").load()
        elif ext == ".pdf":
            res = _load_pdf(path)
            if res is not None:
                # 已按章节聚合：section_path / figure_paths 就绪，直接返回
                return res
            # 降级：PyPDFLoader（纯文本，无结构/无表格图）
            from langchain_community.document_loaders import PyPDFLoader
            docs = PyPDFLoader(path).load()
            out = []
            for d in docs:
                pg = d.metadata.get("page") if hasattr(d, "metadata") else None
                page = (int(pg) + 1) if pg is not None else None
                out.append(RawDoc(
                    text=d.page_content or "",
                    source=path,
                    file_name=os.path.basename(path),
                    page=page,
                ))
            # 降级路径：通用图抽取按页对齐补 figure_paths
            figs_per_page = _extract_figures(path)
            for raw_doc, figs in zip(out, figs_per_page):
                if figs:
                    raw_doc.figure_paths.extend(figs)
            return out
        elif ext in (".html", ".htm"):
            try:
                from langchain_community.document_loaders import BSHTMLLoader
                docs = BSHTMLLoader(path).load()
            except Exception:
                docs = TextLoader(path, encoding="utf-8").load()
        elif ext == ".docx":
            try:
                from langchain_community.document_loaders import Docx2txtLoader
                docs = Docx2txtLoader(path).load()
            except Exception:
                from langchain_community.document_loaders import (
                    UnstructuredWordDocumentLoader,
                )
                docs = UnstructuredWordDocumentLoader(path).load()
        elif ext in (".xlsx", ".xls"):
            res = _load_xlsx(path)
            if res is None:
                print(f"[ingest] 跳过 {path}：缺少 openpyxl 依赖")
                return []
            return res
        elif ext == ".pptx":
            res = _load_pptx(path)
            if res is None:
                print(f"[ingest] 跳过 {path}：缺少 python-pptx 依赖")
                return []
            return res
        else:
            print(f"[ingest] 跳过不支持的格式: {path}")
            return []
    except ImportError as e:
        print(f"[ingest] 缺少依赖，跳过 {path}: {e}")
        return []
    except Exception as e:
        print(f"[ingest] 读取失败 {path}: {e}")
        return []

    out = []
    for d in docs:
        # PyPDFLoader 按页返回；page 转 1-based 人类可读页码（langchain 默认 0-based）
        pg = d.metadata.get("page") if hasattr(d, "metadata") else None
        page = (int(pg) + 1) if pg is not None else None
        out.append(RawDoc(
            text=d.page_content or "",
            source=path,
            file_name=os.path.basename(path),
            page=page,
        ))
    # 通用 PDF 图抽取：每页可能多张图（PyMuPDF + 像素墨迹 + 连通分量）
    # 与语言/排版/有无 caption 无关；表格自然被文字密度过滤；模板背景在像素层自动消失
    figs_per_page = _extract_figures(path)
    for raw_doc, figs in zip(out, figs_per_page):
        if figs:
            raw_doc.figure_paths.extend(figs)
    return out


# --------------------------------------------------------------------------- #
# 通用 PDF 图抽取（PyMuPDF + numpy + scipy，与语言/caption/排版无关）
# --------------------------------------------------------------------------- #
# 算法：
#   1) PyMuPDF 渲染整页 → 像素墨迹
#   2) page.get_text("blocks") 拿文字块 bbox → 文字遮罩（膨胀 4-5 px）
#   3) 图形墨迹 = 整页墨迹 − 文字墨迹
#   4) scipy.ndimage.label 连通分量 → 逐分量判定面积/宽高比/墨迹密度/文字密度
#   5) 通过的真图分量 → 从原图裁剪（保留图内文字标签）→ fig_p{NNN}_{k}.png
#
# 为什么通用：
#   - 无 caption 正则：不依赖页面文本是否含「图N / XX图」，英文 PDF / 无 caption 图照样识别
#   - 模板/背景图层自动消失：get_drawings 返回的伪图元在渲染像素层不存在，就不进 graphic ink
#   - 表格自然被过滤：表格连通分量 bbox 内塞满文字（高密度）→ 丢弃
#   - logo/分隔线/细线：面积阈值与墨迹密度过滤
# --------------------------------------------------------------------------- #
def _extract_figures(path: str) -> List[List[str]]:
    """对 PDF 每一页做通用图抽取；每页返回 0..N 张裁剪图相对路径列表。

    - 缺失 fitz/numpy/scipy/Pillow：返回 []（graceful 降级，pipeline 仍可跑）
    - 非 PDF：返回 []
    - 自动清理旧版 page_*.png（整页渲染）与 fig_pNNN.png（旧单图命名）
    """
    if not path.lower().endswith(".pdf"):
        return []
    try:
        import fitz  # PyMuPDF
        import numpy as np
        from PIL import Image, ImageDraw, ImageFilter
        from scipy import ndimage
    except ImportError as e:
        print(f"[ingest] 跳过 PDF 图抽取 {os.path.basename(path)}：依赖缺失 {e.name}（pip install pymupdf numpy scipy pillow）")
        return []

    try:
        project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        if not os.path.isdir(os.path.join(project_root, "knowledge")):
            project_root = os.getcwd()
        stem = os.path.splitext(os.path.basename(path))[0]
        out_dir = os.path.join(project_root, "assets", "figures", stem)
        os.makedirs(out_dir, exist_ok=True)
        # 清理旧版产物：page_*.png（整页渲染）、fig_pNNN.png（旧单图命名）
        for old in (
            glob.glob(os.path.join(out_dir, "page_*.png"))
            + glob.glob(os.path.join(out_dir, "fig_p[0-9][0-9][0-9].png"))
        ):
            try:
                os.remove(old)
            except OSError:
                pass

        scale = 2.0
        pad = 24
        MIN_AREA_RATIO = 0.01   # 图区面积 / 页面面积
        MAX_INK_DENSITY = 0.12  # 墨迹密度上限（过高=填充色块/水印，非真图）
        MIN_INK_DENSITY = 0.015 # 图区墨迹像素 / bbox 面积（过滤稀疏虚线）
        MAX_LONG_ROWS = 0       # 跨 80% bbox 宽度的长横线行数（>此值=表格，丢弃）；流程图箭头是斜的→不存在长横线；表格分隔线是横的→必然命中

        results: List[List[str]] = []
        with fitz.open(path) as doc:
            for i in range(doc.page_count):
                page_figs: List[str] = []
                try:
                    page = doc[i]
                    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))
                    img = pix.pil_image()
                    W, H = img.size
                    # 文字遮罩
                    blocks = page.get_text("blocks")
                    m = Image.new("L", (W, H), 0)
                    md = ImageDraw.Draw(m)
                    for b in blocks:
                        x0, y0, x1, y1 = b[:4]
                        md.rectangle(
                            [x0*scale-4, y0*scale-4, x1*scale+4, y1*scale+4],
                            fill=255,
                        )
                    m = m.filter(ImageFilter.MaxFilter(5))
                    tm = np.asarray(m) > 128
                    ink = np.asarray(img.convert("L")) < 200
                    graphic = ink & ~tm
                    # 连通分量
                    labeled, ncomp = ndimage.label(graphic)
                    for k in range(1, ncomp + 1):
                        ys, xs = np.where(labeled == k)
                        if len(xs) == 0:
                            continue
                        bx0, bx1 = int(xs.min()), int(xs.max())
                        by0, by1 = int(ys.min()), int(ys.max())
                        bw, bh = bx1 - bx0, by1 - by0
                        if bw <= 0 or bh <= 0:
                            continue
                        if (bw * bh) / (W * H) < MIN_AREA_RATIO:
                            continue
                        asp = bw / bh if bh else 0
                        if asp < 0.1 or asp > 10:
                            continue
                        ink_density = len(xs) / (bw * bh)
                        if ink_density < MIN_INK_DENSITY or ink_density > MAX_INK_DENSITY:
                            continue
                        # 表格判别：长横线（跨 80% bbox 宽度的连续墨迹行）数量
                        sub = np.ascontiguousarray(graphic[by0:by1+1, bx0:bx1+1])
                        long_rows = 0
                        for row in sub:
                            if not row.any():
                                continue
                            diffs = np.diff(np.r_[0, row.astype(np.int8), 0])
                            starts = np.where(diffs == 1)[0]
                            ends = np.where(diffs == -1)[0]
                            max_run = max((e - s) for s, e in zip(starts, ends))
                            if max_run >= 0.8 * bw:
                                long_rows += 1
                        if long_rows > MAX_LONG_ROWS:
                            continue
                        # 裁剪原图（保留图内文字标签）
                        crop = img.crop((
                            max(0, bx0 - pad),
                            max(0, by0 - pad),
                            min(W, bx1 + pad),
                            min(H, by1 + pad),
                        ))
                        fname = f"fig_p{i+1:03d}_{len(page_figs)+1}.png"
                        crop.save(os.path.join(out_dir, fname), format="PNG", optimize=True)
                        page_figs.append(f"assets/figures/{stem}/{fname}")

                    # 兜底：页面含图标题锚点但连通分量没抽出真图时（常见流程图/时序图
                    # 被长横线过滤误伤），直接渲染整页作为 fallback 图，确保 LLM 能看到。
                    if not page_figs and FIGURE_ANCHOR_RE.search(page.get_text() or ""):
                        try:
                            fname = f"page_{i+1:03d}.png"
                            page.get_pixmap(matrix=fitz.Matrix(scale, scale)).save(
                                os.path.join(out_dir, fname))
                            page_figs.append(f"assets/figures/{stem}/{fname}")
                        except Exception as fe:
                            print(f"[ingest] PDF 第 {i+1} 页整页渲染兜底失败: {fe}")
                except Exception as e:
                    print(f"[ingest] PDF 第 {i+1} 页图抽取失败: {e}")
                results.append(page_figs)
        return results
    except Exception as e:
        print(f"[ingest] PDF 图抽取异常 {os.path.basename(path)}: {e}")
        return []
